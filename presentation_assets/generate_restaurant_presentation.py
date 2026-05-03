from __future__ import annotations

import sqlite3
import subprocess
from pathlib import Path

import qrcode
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Pt as DocPt
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path("/Users/saidbekdavranbekov/Documents/Диплом ресторан")
ASSETS_DIR = BASE_DIR / "presentation_assets"
OUTPUT_PPTX = BASE_DIR / "Презентация_ДП_Гастрономия.pptx"
OUTPUT_PPTX_CLEAN = BASE_DIR / "Презентация_ДП_Гастрономия_без_кругов.pptx"
OUTPUT_PPTX_SOFT = BASE_DIR / "Презентация_ДП_Гастрономия_с_кругами.pptx"
OUTPUT_DOCX = ASSETS_DIR / "Тест-кейсы_Гастрономия.docx"
OUTPUT_TESTCASE_PNG = ASSETS_DIR / "test_cases_word_style.png"
OUTPUT_QR = ASSETS_DIR / "qr_demo.png"
TEMPLATE_PPTX = Path("/var/folders/kf/f94mvc650dxd8sdfw26lfs9m0000gn/T/преза.pptx")
ROUNDED_DIR = ASSETS_DIR / "rounded"

HOME_SCREENSHOT = ASSETS_DIR / "home.png"
MENU_SCREENSHOT = ASSETS_DIR / "menu.png"
ADMIN_SCREENSHOT = ASSETS_DIR / "admin.png"

FONT_REGULAR = "/System/Library/Fonts/Supplemental/Arial.ttf"
FONT_BOLD = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"

DARK_BG = RGBColor(16, 22, 35)
PANEL = RGBColor(25, 35, 53)
PANEL_ALT = RGBColor(31, 43, 64)
ACCENT = RGBColor(255, 122, 0)
ACCENT_ALT = RGBColor(255, 61, 84)
TEXT_MAIN = RGBColor(248, 249, 251)
TEXT_MUTED = RGBColor(178, 184, 196)
LINE = RGBColor(62, 78, 110)
SUCCESS = RGBColor(68, 201, 126)
WARNING = RGBColor(255, 184, 77)
ERROR = RGBColor(255, 99, 99)


def get_local_ip() -> str:
    for iface in ("en0", "en1"):
        try:
            return (
                subprocess.check_output(["ipconfig", "getifaddr", iface], text=True)
                .strip()
            )
        except subprocess.CalledProcessError:
            continue
    return "127.0.0.1"


def get_db_counts() -> dict[str, int]:
    db_path = BASE_DIR / "Документация/back/restaurant.db"
    conn = sqlite3.connect(db_path)
    cur = conn.cursor()
    result = {}
    for table in ("menu_items", "reviews", "gallery_images", "orders", "users"):
        cur.execute(f"SELECT COUNT(*) FROM {table}")
        result[table] = cur.fetchone()[0]
    conn.close()
    return result


def pil_font(size: int, bold: bool = False):
    return ImageFont.truetype(FONT_BOLD if bold else FONT_REGULAR, size=size)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> str:
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        test = f"{current} {word}".strip()
        width = draw.textbbox((0, 0), test, font=font)[2]
        if width <= max_width or not current:
            current = test
        else:
            lines.append(current)
            current = word
    if current:
        lines.append(current)
    return "\n".join(lines)


def generate_qr_image(url: str) -> None:
    qr = qrcode.QRCode(border=1, box_size=10)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white")
    img.save(OUTPUT_QR)


def generate_test_cases_docx() -> None:
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Cm(1.5)
    sec.bottom_margin = Cm(1.5)
    sec.left_margin = Cm(1.5)
    sec.right_margin = Cm(1.5)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("Фрагмент таблицы тест-кейсов\nВеб-приложение ресторана «Гастрономия»")
    run.bold = True
    run.font.size = DocPt(14)

    table = doc.add_table(rows=3, cols=6)
    table.style = "Table Grid"
    headers = ["ID", "Проверка", "Шаги", "Ожидаемый результат", "Факт", "Статус"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        for p in cell.paragraphs:
            for r in p.runs:
                r.bold = True
                r.font.size = DocPt(10)

    rows = [
        [
            "TC-01",
            "Бронирование стола с валидным номером",
            "Открыть форму → ввести +7 (999) 123-45-67 → нажать «Подтвердить»",
            "Появляется окно успешного бронирования, запрос уходит на /api/table-reservations",
            "Форма отправлена, пользователь видит окно успеха",
            "Пройден",
        ],
        [
            "TC-02",
            "Редактирование кода города в маске телефона",
            "Открыть форму → ввести +7 (923) → нажать Backspace/Delete возле )",
            "Цифры кода города удаляются, номер можно исправить",
            "Курсор останавливался на скобке, цифры нельзя было удалить",
            "Провален",
        ],
    ]
    for row_idx, data in enumerate(rows, start=1):
        for col_idx, value in enumerate(data):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            for p in cell.paragraphs:
                for r in p.runs:
                    r.font.size = DocPt(9)
    doc.save(OUTPUT_DOCX)


def generate_test_cases_image() -> None:
    width, height = 1800, 1050
    img = Image.new("RGB", (width, height), "#e8eaed")
    draw = ImageDraw.Draw(img)

    page = (80, 70, width - 80, height - 70)
    draw.rounded_rectangle(page, radius=20, fill="white", outline="#d0d7de", width=2)
    draw.rectangle((page[0], page[1], page[2], page[1] + 54), fill="#f3f4f6")
    draw.ellipse((110, 88, 126, 104), fill="#ef4444")
    draw.ellipse((136, 88, 152, 104), fill="#f59e0b")
    draw.ellipse((162, 88, 178, 104), fill="#10b981")
    draw.text((page[0] + 620, 85), "Microsoft Word — Тест-кейсы_Гастрономия.docx", fill="#374151", font=pil_font(22, bold=True))

    draw.text((page[0] + 60, page[1] + 90), "Фрагмент таблицы тест-кейсов", fill="#111827", font=pil_font(30, bold=True))
    draw.text((page[0] + 60, page[1] + 130), "Для презентации выбраны один пройденный и один проваленный кейс.", fill="#6b7280", font=pil_font(20))

    headers = ["ID", "Проверка", "Шаги", "Ожидаемо", "Фактически", "Статус"]
    rows = [
        [
            "TC-01",
            "Бронирование с валидным номером",
            "Открыть форму\nВвести +7 (999) 123-45-67\nНажать «Подтвердить»",
            "Окно успеха\nPOST /api/table-reservations",
            "Успешно выполнено",
            "Пройден",
        ],
        [
            "TC-02",
            "Редактирование кода города",
            "Открыть форму\nВвести +7 (923)\nНажать Backspace/Delete возле )",
            "Цифры в скобках удаляются,\nномер корректируется",
            "Курсор стопорится на скобке,\nцифры не удаляются",
            "Провален",
        ],
    ]

    table_left = page[0] + 60
    table_top = page[1] + 190
    col_widths = [120, 250, 310, 290, 290, 150]
    row_height = 190

    x = table_left
    for idx, header in enumerate(headers):
        fill = "#eef2ff" if idx < 5 else "#fff7ed"
        draw.rectangle((x, table_top, x + col_widths[idx], table_top + 70), fill=fill, outline="#cbd5e1", width=2)
        wrapped = wrap_text(draw, header, pil_font(20, bold=True), col_widths[idx] - 20)
        draw.multiline_text((x + 10, table_top + 15), wrapped, fill="#111827", font=pil_font(20, bold=True), spacing=4)
        x += col_widths[idx]

    for row_idx, row in enumerate(rows):
        y1 = table_top + 70 + row_idx * row_height
        y2 = y1 + row_height
        x = table_left
        for col_idx, value in enumerate(row):
            status_fill = "#ecfdf5" if value == "Пройден" else "#fef2f2" if value == "Провален" else "white"
            draw.rectangle((x, y1, x + col_widths[col_idx], y2), fill=status_fill, outline="#d1d5db", width=2)
            wrapped = wrap_text(draw, value, pil_font(18, bold=(col_idx == 0 or col_idx == 5)), col_widths[col_idx] - 18)
            fill = "#065f46" if value == "Пройден" else "#991b1b" if value == "Провален" else "#111827"
            draw.multiline_text((x + 9, y1 + 10), wrapped, fill=fill, font=pil_font(18, bold=(col_idx == 0 or col_idx == 5)), spacing=5)
            x += col_widths[col_idx]

    img.save(OUTPUT_TESTCASE_PNG)


def clear_slides(prs: Presentation) -> None:
    for i in range(len(prs.slides) - 1, -1, -1):
        r_id = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[i]


def add_bg(slide, variant: str = "soft") -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    if variant == "clean":
        return
    for args in [
        (-1.15, 6.58, 2.05, 1.65, ACCENT, 0.9),
        (10.85, -0.95, 2.45, 2.1, ACCENT_ALT, 0.9),
        (9.75, 6.58, 2.55, 1.7, ACCENT, 0.92),
    ]:
        glow = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(args[0]),
            Inches(args[1]),
            Inches(args[2]),
            Inches(args[3]),
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = args[4]
        glow.fill.transparency = args[5]
        glow.line.fill.background()


def add_footer(slide, page_num: int) -> None:
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(4.8), Inches(0.25))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Дипломный проект • Ресторан «Гастрономия»"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED

    num = slide.shapes.add_textbox(Inches(12.1), Inches(7.0), Inches(0.6), Inches(0.3))
    tf2 = num.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    p2.text = f"{page_num:02d}"
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = TEXT_MUTED


def add_title(slide, title: str, subtitle: str | None = None, badge: str | None = None) -> None:
    if badge:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.42), Inches(1.85), Inches(0.38))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = badge
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = TEXT_MAIN

    tb = slide.shapes.add_textbox(Inches(0.62), Inches(0.9), Inches(11.5), Inches(0.65))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.63), Inches(2.2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.65), Inches(1.78), Inches(11.2), Inches(0.5))
        tf2 = st.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MUTED


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=TEXT_MAIN,
    bold=False,
    align=PP_ALIGN.LEFT,
    fit=False,
    valign=MSO_ANCHOR.TOP,
    margin_left=0.04,
    margin_right=0.04,
    margin_top=0.02,
    margin_bottom=0.02,
    line_spacing=1.0,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(margin_right)
    tf.margin_top = Inches(margin_top)
    tf.margin_bottom = Inches(margin_bottom)

    if fit:
        tf.text = str(text)
        try:
            tf.fit_text(
                font_family="Arial",
                max_size=font_size,
                bold=bold,
                italic=False,
                font_file=FONT_BOLD if bold else FONT_REGULAR,
            )
            for p in tf.paragraphs:
                p.alignment = align
                p.line_spacing = line_spacing
                p.space_after = Pt(0)
                for run in p.runs:
                    run.font.color.rgb = color
                    run.font.bold = bold
            return tb
        except (TypeError, ValueError):
            tf.clear()

    if not fit or not tf.text:
        for idx, line in enumerate(str(text).split("\n")):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = align
            p.line_spacing = line_spacing
            p.space_after = Pt(0)
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
    return tb


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title,
    body,
    fill=PANEL,
    accent=ACCENT,
    body_size=14,
    title_size=16,
    fit_body=False,
    title_top=0.12,
    body_top=0.52,
    body_bottom=0.12,
):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.2)

    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.08), height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    add_textbox(
        slide,
        left + Inches(0.18),
        top + Inches(title_top),
        width - Inches(0.28),
        Inches(0.3),
        title,
        font_size=title_size,
        bold=True,
    )
    add_textbox(
        slide,
        left + Inches(0.18),
        top + Inches(body_top),
        width - Inches(0.3),
        height - Inches(body_top + body_bottom),
        body,
        font_size=body_size,
        color=TEXT_MUTED,
        fit=fit_body,
        margin_left=0.02,
        margin_right=0.03,
        margin_top=0.01,
        margin_bottom=0.02,
        line_spacing=1.02,
    )
    return shape


def add_metric(slide, left, top, width, height, number, label, number_size=22, label_size=11):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL_ALT
    shape.line.color.rgb = LINE
    add_textbox(slide, left, top + Inches(0.1), width, Inches(0.42), number, font_size=number_size, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(
        slide,
        left + Inches(0.08),
        top + Inches(0.64),
        width - Inches(0.16),
        height - Inches(0.7),
        label,
        font_size=label_size,
        color=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
        fit=True,
        margin_left=0.02,
        margin_right=0.02,
    )


def prepare_rounded_cover(image_path: Path, cache_key: str, target_width: int, target_height: int, radius: int = 42) -> Path:
    ROUNDED_DIR.mkdir(parents=True, exist_ok=True)
    output_path = ROUNDED_DIR / f"{cache_key}_{target_width}x{target_height}_r{radius}.png"
    if output_path.exists():
        return output_path

    with Image.open(image_path) as img:
        source = img.convert("RGBA")
        src_w, src_h = source.size
        target_ratio = target_width / target_height
        src_ratio = src_w / src_h

        if src_ratio > target_ratio:
            crop_w = int(src_h * target_ratio)
            left = (src_w - crop_w) // 2
            crop_box = (left, 0, left + crop_w, src_h)
        else:
            crop_h = int(src_w / target_ratio)
            top = (src_h - crop_h) // 2
            crop_box = (0, top, src_w, top + crop_h)

        cropped = source.crop(crop_box).resize((target_width, target_height), Image.Resampling.LANCZOS)
        mask = Image.new("L", (target_width, target_height), 0)
        ImageDraw.Draw(mask).rounded_rectangle((0, 0, target_width, target_height), radius=radius, fill=255)
        rounded = Image.new("RGBA", (target_width, target_height), (0, 0, 0, 0))
        rounded.paste(cropped, (0, 0), mask)
        rounded.save(output_path)

    return output_path


def add_picture_cover(slide, image_path: Path, left, top, width, height, rounded=False, cache_key: str | None = None):
    source_path = image_path
    if rounded:
        source_path = prepare_rounded_cover(
            image_path,
            cache_key or image_path.stem,
            max(960, int(width.inches * 320)),
            max(620, int(height.inches * 320)),
        )
    pic = slide.shapes.add_picture(str(source_path), left, top, width=width, height=height)
    if rounded:
        return pic
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    image_ratio = img_w / img_h
    shape_ratio = width / height
    if image_ratio > shape_ratio:
        visible_width = img_h * shape_ratio
        crop = (img_w - visible_width) / img_w / 2
        pic.crop_left = crop
        pic.crop_right = crop
    else:
        visible_height = img_w / shape_ratio
        crop = (img_h - visible_height) / img_h / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def add_connector(slide, x1, y1, x2, y2, color=ACCENT):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    return conn


def add_entity(slide, left, top, width, height, name, fields):
    outer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    outer.fill.solid()
    outer.fill.fore_color.rgb = PANEL
    outer.line.color.rgb = LINE
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, Inches(0.42))
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT
    header.line.fill.background()
    add_textbox(slide, left, top + Inches(0.03), width, Inches(0.25), name, font_size=15, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.08), top + Inches(0.5), width - Inches(0.15), height - Inches(0.58), "\n".join(fields), font_size=11, color=TEXT_MUTED)


def build_presentation(output_path: Path = OUTPUT_PPTX, bg_variant: str = "clean"):
    ASSETS_DIR.mkdir(exist_ok=True)
    ip = get_local_ip()
    counts = get_db_counts()
    qr_url = f"http://{ip}:3000"
    generate_qr_image(qr_url)
    generate_test_cases_docx()
    generate_test_cases_image()

    prs = Presentation(str(TEMPLATE_PPTX))
    clear_slides(prs)
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.66), Inches(0.62), Inches(2.05), Inches(0.4))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    add_textbox(slide, Inches(0.75), Inches(0.69), Inches(1.9), Inches(0.2), "ДИПЛОМНЫЙ ПРОЕКТ", font_size=11, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.68), Inches(1.28), Inches(8.6), Inches(1.15), "Разработка веб-приложения\nдля ресторана «Гастрономия»", font_size=28, bold=True)
    add_textbox(slide, Inches(0.7), Inches(2.5), Inches(8.3), Inches(0.5), "SPA на Vue 3 + FastAPI backend + SQLite + Docker-деплой", font_size=16, color=TEXT_MUTED)
    add_card(
        slide,
        Inches(0.72),
        Inches(5.35),
        Inches(3.45),
        Inches(1.18),
        "Студент",
        "Давранбеков С.Б.\nНаправление: информационные\nсистемы",
        body_size=13,
        fit_body=True,
    )
    add_card(
        slide,
        Inches(4.45),
        Inches(5.35),
        Inches(3.55),
        Inches(1.18),
        "Предметная область",
        "Автоматизация презентации ресторана,\nонлайн-заказов и бронирования",
        body_size=13,
        fit_body=True,
    )
    add_card(
        slide,
        Inches(8.25),
        Inches(5.35),
        Inches(4.02),
        Inches(1.18),
        "Результат",
        "Публичный сайт, админ-панель,\nREST API и Docker-сборка на 3000 порту",
        body_size=13,
        fit_body=True,
    )

    # 2
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Цель дипломного проекта", "Цель и ожидаемый эффект от внедрения цифрового сервиса", "02")
    add_card(
        slide,
        Inches(0.7),
        Inches(2.2),
        Inches(7.2),
        Inches(2.3),
        "Цель ДП",
        "Разработать веб-приложение ресторана «Гастрономия», которое объединяет презентацию бренда, онлайн-бронирование, оформление заказа и административное управление контентом в едином интерфейсе.",
        body_size=17,
        fit_body=True,
    )
    add_card(
        slide,
        Inches(0.7),
        Inches(4.7),
        Inches(7.2),
        Inches(1.45),
        "Ожидаемый эффект",
        "Повышение конверсии из посетителя в заявку, сокращение ручной обработки заказов и ускорение обновления контента через админ-панель.",
        body_size=15,
        fit_body=True,
    )
    add_metric(slide, Inches(8.2), Inches(2.2), Inches(1.35), Inches(1.2), str(counts["menu_items"]), "позиций меню\nв БД", label_size=10)
    add_metric(slide, Inches(9.72), Inches(2.2), Inches(1.35), Inches(1.2), str(counts["reviews"]), "отзыва\nв системе", label_size=10)
    add_metric(slide, Inches(11.24), Inches(2.2), Inches(1.35), Inches(1.2), str(counts["gallery_images"]), "изображений\nгалереи", label_size=10)
    add_card(
        slide,
        Inches(8.15),
        Inches(3.75),
        Inches(4.15),
        Inches(2.12),
        "Ключевые показатели",
        "• 5 маршрутов SPA\n• 2 роли пользователей\n• REST API для меню, отзывов, заказов и галереи\n• Docker-развертывание для демонстрации",
        body_size=15,
        fit_body=True,
    )
    add_footer(slide, 2)

    # 3
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Сбор требований", "Фрагмент интервью с заказчиком и ожидания целевой аудитории", "03")
    add_card(slide, Inches(0.7), Inches(2.2), Inches(5.8), Inches(3.9), "Требования заказчика", "1. Современный визуальный сайт в фирменном стиле ресторана.\n2. Публичное меню с ценами, составом и фотографиями.\n3. Формы онлайн-заказа и бронирования стола без звонка.\n4. Раздел отзывов и галерея для повышения доверия.\n5. Админ-панель для самостоятельного обновления данных.", body_size=15, fit_body=True)
    add_card(slide, Inches(6.8), Inches(2.2), Inches(5.55), Inches(3.9), "Ожидания целевой аудитории", "1. Быстро понять ассортимент и стоимость блюд.\n2. Удобно пользоваться сайтом со смартфона.\n3. Сразу видеть фото интерьера и блюд.\n4. Забронировать стол за 1–2 действия.\n5. Видеть актуальные отзывы и подтверждение после отправки формы.", fill=PANEL_ALT, body_size=15, fit_body=True)
    add_footer(slide, 3)

    # 4
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "На основании требований формируем задачи", "Декомпозиция целей в конкретные задачи разработки", "04")
    tasks = [
        ("01", "Спроектировать структуру SPA и пользовательские маршруты."),
        ("02", "Разработать REST API для меню, отзывов, галереи, заказов и бронирования."),
        ("03", "Реализовать админ-панель с JWT-аутентификацией и разграничением прав."),
        ("04", "Спроектировать БД SQLite и инициализировать ее стартовыми данными."),
        ("05", "Провести тестирование интерфейса, API и критических пользовательских сценариев."),
        ("06", "Подготовить контейнеризацию проекта и удобный демонстрационный запуск."),
    ]
    x_positions = [0.7, 6.65]
    y = 2.1
    idx = 0
    for row in range(3):
        for col in range(2):
            num, text = tasks[idx]
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x_positions[col]), Inches(y), Inches(5.6), Inches(1.25))
            card.fill.solid()
            card.fill.fore_color.rgb = PANEL
            card.line.color.rgb = LINE
            circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x_positions[col] + 0.18), Inches(y + 0.18), Inches(0.7), Inches(0.7))
            circle.fill.solid()
            circle.fill.fore_color.rgb = ACCENT
            circle.line.fill.background()
            add_textbox(slide, Inches(x_positions[col] + 0.18), Inches(y + 0.35), Inches(0.7), Inches(0.2), num, font_size=14, bold=True, align=PP_ALIGN.CENTER)
            add_textbox(slide, Inches(x_positions[col] + 1.0), Inches(y + 0.18), Inches(4.35), Inches(0.75), text, font_size=15, fit=True)
            idx += 1
        y += 1.42
    add_footer(slide, 4)

    # 5
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Диаграмма связей (mindmap)", "Основные подсистемы и направления разработки", "05")
    center_left, center_top = Inches(4.2), Inches(3.25)
    center_width, center_height = Inches(4.6), Inches(1.05)
    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, center_left, center_top, center_width, center_height)
    center.fill.solid()
    center.fill.fore_color.rgb = ACCENT
    center.line.fill.background()
    add_textbox(
        slide,
        center_left + Inches(0.15),
        center_top + Inches(0.16),
        center_width - Inches(0.3),
        Inches(0.6),
        "Веб-приложение ресторана\n«Гастрономия»",
        font_size=20,
        bold=True,
        align=PP_ALIGN.CENTER,
        fit=True,
    )
    nodes = [
        (Inches(0.78), Inches(2.55), Inches(3.1), Inches(1.08), "Frontend", "Vue 3\nVue Router\nVuex / Axios", "left"),
        (Inches(0.78), Inches(4.95), Inches(3.1), Inches(1.08), "Коммерция", "Корзина\nЗаказ\nБронирование", "left"),
        (Inches(4.45), Inches(2.25), Inches(3.1), Inches(1.0), "Контент", "Меню\nОтзывы\nГалерея", "top"),
        (Inches(9.22), Inches(2.55), Inches(3.05), Inches(1.08), "Администрирование", "JWT-вход\nCRUD\nСтатусы заказов", "right"),
        (Inches(9.22), Inches(4.95), Inches(3.05), Inches(1.08), "Backend / Data", "FastAPI\nSQLAlchemy\nSQLite", "right"),
        (Inches(4.45), Inches(5.85), Inches(3.1), Inches(1.0), "DevOps", "Docker Compose\nNginx\nРазвертывание", "bottom"),
    ]
    for left, top, width, height, title, body, anchor in nodes:
        if anchor == "left":
            add_connector(slide, center_left, center_top + center_height / 2, left + width, top + height / 2)
        elif anchor == "right":
            add_connector(slide, center_left + center_width, center_top + center_height / 2, left, top + height / 2)
        elif anchor == "top":
            add_connector(slide, center_left + center_width / 2, center_top, left + width / 2, top + height)
        else:
            add_connector(slide, center_left + center_width / 2, center_top + center_height, left + width / 2, top)
    for left, top, width, height, title, body, anchor in nodes:
        add_card(slide, left, top, width, height, title, body, body_size=13, fit_body=True)
    add_footer(slide, 5)

    # 6
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Сценарии использования проектируемого ПП", "Основные use-case для гостя и администратора", "06")
    boundary = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(2.55), Inches(2.05), Inches(8.2), Inches(4.65))
    boundary.fill.solid()
    boundary.fill.fore_color.rgb = PANEL
    boundary.fill.transparency = 0.08
    boundary.line.color.rgb = LINE
    add_textbox(slide, Inches(5.0), Inches(2.12), Inches(3.2), Inches(0.28), "Система «Гастрономия»", font_size=16, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, Inches(0.48), Inches(3.08), Inches(1.82), Inches(1.3), "Гость", "Публичный\nпользователь", body_size=11, fit_body=True, body_top=0.48, body_bottom=0.08)
    add_card(slide, Inches(10.78), Inches(3.08), Inches(1.94), Inches(1.3), "Админ", "Авторизованный\nсотрудник", body_size=11, fill=PANEL_ALT, fit_body=True, body_top=0.48, body_bottom=0.08)
    add_textbox(slide, Inches(3.05), Inches(2.42), Inches(2.55), Inches(0.2), "Публичные сценарии", font_size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.0), Inches(2.42), Inches(2.75), Inches(0.2), "Административные сценарии", font_size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    guest_cases = [
        (3.0, 2.75, "Просмотреть меню"),
        (3.0, 3.68, "Оставить отзыв"),
        (3.0, 4.61, "Оформить заказ"),
        (3.0, 5.54, "Забронировать стол"),
    ]
    admin_cases = [
        (7.0, 2.75, "Авторизоваться"),
        (7.0, 3.68, "Управлять меню"),
        (7.0, 4.61, "Управлять отзывами"),
        (7.0, 5.54, "Управлять заказами"),
    ]
    for x, y, text in guest_cases + admin_cases:
        oval = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(2.7), Inches(0.62))
        oval.fill.solid()
        oval.fill.fore_color.rgb = PANEL_ALT
        oval.line.color.rgb = ACCENT
        add_textbox(slide, Inches(x + 0.12), Inches(y + 0.17), Inches(2.46), Inches(0.18), text, font_size=13, align=PP_ALIGN.CENTER, fit=True)
    for x, y, _ in guest_cases:
        add_connector(slide, Inches(2.3), Inches(3.73), Inches(x), Inches(y + 0.31))
    for x, y, _ in admin_cases:
        add_connector(slide, Inches(10.78), Inches(3.73), Inches(x + 2.7), Inches(y + 0.31))
    add_footer(slide, 6)

    # 7
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Прототипы / макеты", "Реальные экраны реализованного интерфейса", "07")
    panels = [
        (HOME_SCREENSHOT, "Главная страница", "презентация ресторана, CTA и блоки контента"),
        (MENU_SCREENSHOT, "Каталог меню", "фильтрация категорий и добавление в корзину"),
        (ADMIN_SCREENSHOT, "Админ-панель", "точка входа для управления контентом"),
    ]
    x = 0.68
    for path, title, desc in panels:
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.1), Inches(4.0), Inches(3.7))
        frame.fill.solid()
        frame.fill.fore_color.rgb = PANEL
        frame.line.color.rgb = LINE
        if path.exists():
            add_picture_cover(slide, path, Inches(x + 0.08), Inches(2.18), Inches(3.84), Inches(2.5), rounded=True, cache_key=f"{path.stem}_slide7")
        add_textbox(slide, Inches(x + 0.12), Inches(4.82), Inches(3.7), Inches(0.25), title, font_size=14, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.16), Inches(5.1), Inches(3.65), Inches(0.45), desc, font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        x += 4.18
    add_footer(slide, 7)

    # 8
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Архитектура ПП", "Логическая схема развертывания и взаимодействия компонентов", "08")
    layers = [
        (0.75, 2.45, 2.35, "Клиентский слой", "Браузер\nVue 3 SPA\nVue Router / Vuex"),
        (3.35, 2.45, 2.45, "Web-слой", "Nginx\nстатический frontend\nпрокси /api и /images"),
        (6.1, 2.45, 2.45, "Сервер", "FastAPI\nREST API\nJWT, CORS, upload"),
        (8.9, 2.45, 3.1, "Хранение данных", "SQLite\nкаталог images/\nseed-данные и backup"),
    ]
    for x, y, w, title, body in layers:
        add_card(slide, Inches(x), Inches(y), Inches(w), Inches(1.82), title, body, body_size=14, fit_body=True)
    add_connector(slide, Inches(3.1), Inches(3.36), Inches(3.35), Inches(3.36))
    add_connector(slide, Inches(5.8), Inches(3.36), Inches(6.1), Inches(3.36))
    add_connector(slide, Inches(8.55), Inches(3.36), Inches(8.9), Inches(3.36))
    add_card(slide, Inches(0.85), Inches(5.0), Inches(3.45), Inches(1.45), "Интеграции фронта", "GET/POST через Axios:\n/api/menu, /api/reviews, /api/orders,\n/api/table-reservations", body_size=13, fit_body=True)
    add_card(slide, Inches(4.55), Inches(5.0), Inches(3.2), Inches(1.45), "Безопасность", "JWT для админки,\nвалидация входных данных,\nограничение upload по MIME-типу", body_size=13, fill=PANEL_ALT, fit_body=True)
    add_card(slide, Inches(8.0), Inches(5.0), Inches(4.25), Inches(1.45), "Инфраструктура", "Docker Compose объединяет frontend и backend,\nа демонстрация публикуется на порту 3000.", body_size=13, fit_body=True)
    add_footer(slide, 8)

    # 9
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Структура базы данных (ERD)", "Физическая модель данных проекта на SQLite", "09")
    add_entity(slide, Inches(0.7), Inches(2.1), Inches(2.45), Inches(1.85), "menu_items", ["PK id", "name", "category", "price / weight", "composition", "image", "featured"])
    add_entity(slide, Inches(3.45), Inches(2.1), Inches(2.35), Inches(1.85), "reviews", ["PK id", "name / email", "rating", "text", "answer", "created_at"])
    add_entity(slide, Inches(6.05), Inches(2.1), Inches(2.45), Inches(1.85), "gallery_images", ["PK id", "src", "alt", "description", "category", "created_at"])
    add_entity(slide, Inches(8.8), Inches(2.1), Inches(2.8), Inches(1.85), "orders", ["PK id", "items (JSON)", "total", "customer_phone", "customer_email", "status", "created_at"])
    add_entity(slide, Inches(4.5), Inches(4.4), Inches(3.0), Inches(1.65), "users", ["PK id", "username", "email", "hashed_password", "is_active", "is_admin"])
    add_connector(slide, Inches(6.0), Inches(4.4), Inches(6.0), Inches(3.95), color=TEXT_MUTED)
    add_textbox(slide, Inches(0.9), Inches(6.25), Inches(11.4), Inches(0.45), "Примечание: для таблицы orders состав заказа хранится как JSON-строка в поле items, а связь пользователя с контентом реализуется логически через защищенные API-эндпоинты.", font_size=12, color=TEXT_MUTED)
    add_footer(slide, 9)

    # 10
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Разработка ПП: инструментальные и программные средства", "Технологический стек проекта", "10")
    add_card(slide, Inches(0.75), Inches(2.15), Inches(3.75), Inches(3.25), "Frontend", "• Vue 3\n• Vite\n• Tailwind CSS\n• Vue Router\n• Vuex\n• Axios", body_size=16)
    add_card(slide, Inches(4.8), Inches(2.15), Inches(3.75), Inches(3.25), "Backend", "• FastAPI\n• SQLAlchemy\n• SQLite\n• Pydantic schemas\n• JWT + passlib\n• Upload image API", body_size=16, fill=PANEL_ALT)
    add_card(slide, Inches(8.85), Inches(2.15), Inches(3.45), Inches(3.25), "Инструменты и DevOps", "• VS Code\n• Git / GitHub\n• Docker Compose\n• Nginx\n• Chrome DevTools\n• python-pptx для генерации презентации", body_size=15)
    add_footer(slide, 10)

    # 11
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Календарный план разработки", "Этапы выполнения дипломного проекта по неделям", "11")
    weeks = ["1", "2", "3", "4", "5", "6", "7", "8"]
    stages = [
        ("Сбор требований", 1, 1, ACCENT),
        ("Проектирование UI/UX", 1, 2, WARNING),
        ("Frontend SPA", 2, 4, RGBColor(86, 152, 255)),
        ("Backend + БД", 3, 5, RGBColor(129, 140, 248)),
        ("Интеграция API", 5, 6, RGBColor(45, 212, 191)),
        ("Тестирование и bugfix", 6, 7, SUCCESS),
        ("Docker и подготовка к защите", 7, 8, ACCENT_ALT),
    ]
    start_x, start_y = 0.78, 2.25
    col_w, row_h = 0.78, 0.52
    add_textbox(slide, Inches(start_x), Inches(start_y), Inches(2.1), Inches(0.25), "Этап", font_size=13, bold=True)
    for idx, week in enumerate(weeks):
        header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(start_x + 2.25 + idx * col_w), Inches(start_y - 0.03), Inches(0.66), Inches(0.34))
        header.fill.solid()
        header.fill.fore_color.rgb = PANEL_ALT
        header.line.color.rgb = LINE
        add_textbox(slide, Inches(start_x + 2.25 + idx * col_w), Inches(start_y + 0.04), Inches(0.66), Inches(0.15), week, font_size=12, bold=True, align=PP_ALIGN.CENTER)
    y = start_y + 0.45
    for stage, w1, w2, color in stages:
        add_textbox(slide, Inches(start_x), Inches(y + 0.06), Inches(2.0), Inches(0.2), stage, font_size=12, color=TEXT_MAIN)
        for idx in range(8):
            cell = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(start_x + 2.25 + idx * col_w), Inches(y), Inches(0.66), Inches(0.34))
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL
            cell.line.color.rgb = LINE
        bar_left = start_x + 2.25 + (w1 - 1) * col_w
        bar_width = (w2 - w1 + 1) * col_w - 0.02
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(bar_left), Inches(y + 0.02), Inches(bar_width), Inches(0.3))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        y += row_h
    add_footer(slide, 11)

    # 12
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Доска задач", "Kanban-представление состояния работ по проекту", "12")
    columns = [
        ("Backlog", ["Аналитика посещаемости", "Email/SMS уведомления", "Онлайн-оплата"]),
        ("In Progress", ["Сохранение бронирований в БД", "Расширение формы бронирования"]),
        ("Review", ["UX-сообщения об ошибках", "Документация по развертыванию"]),
        ("Done", ["Маршруты SPA", "Меню / отзывы / галерея", "Корзина и заказ", "JWT-админка", "Docker-развертывание"]),
    ]
    colors = [TEXT_MUTED, WARNING, RGBColor(103, 149, 255), SUCCESS]
    x = 0.65
    for idx, (title, items) in enumerate(columns):
        col = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.95), Inches(4.6))
        col.fill.solid()
        col.fill.fore_color.rgb = PANEL
        col.line.color.rgb = LINE
        head = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.12), Inches(2.12), Inches(2.71), Inches(0.45))
        head.fill.solid()
        head.fill.fore_color.rgb = PANEL_ALT
        head.line.fill.background()
        add_textbox(slide, Inches(x + 0.12), Inches(2.22), Inches(2.71), Inches(0.18), title, font_size=14, bold=True, align=PP_ALIGN.CENTER)
        y = 2.75
        for item in items:
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.18), Inches(y), Inches(2.6), Inches(0.58))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(34, 46, 69)
            card.line.color.rgb = colors[idx]
            add_textbox(slide, Inches(x + 0.28), Inches(y + 0.11), Inches(2.35), Inches(0.25), item, font_size=11)
            y += 0.72
        x += 3.14
    add_footer(slide, 12)

    # 13
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Тест-план / тест-кейсы и bug-report", "Для презентации выбраны один пройденный и один проваленный кейс", "13")
    add_card(slide, Inches(0.72), Inches(2.0), Inches(7.2), Inches(4.45), "Фрагмент тест-кейсов", "На изображении ниже представлен фрагмент таблицы, оформленной в Word-стиле: один кейс успешно пройден, один — провален.", body_size=13)
    if OUTPUT_TESTCASE_PNG.exists():
        add_picture_cover(slide, OUTPUT_TESTCASE_PNG, Inches(0.88), Inches(2.62), Inches(6.9), Inches(3.55))
    add_card(slide, Inches(8.15), Inches(2.0), Inches(4.1), Inches(4.45), "Bug-report: BUG-02", "Название: нельзя удалить цифры кода города, если курсор упирается в )\n\nШаги:\n1. Открыть форму бронирования.\n2. Ввести +7 (923)\n3. Нажать Backspace/Delete рядом со скобкой.\n\nОжидаемо: цифры удаляются, номер корректируется.\nФактически: курсор стопорится на символе ')'.\n\nПриоритет: Medium\nСтатус: Исправлено", fill=PANEL_ALT, body_size=13)
    add_footer(slide, 13)

    # 14
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Информационная безопасность", "Базовые меры защиты, реализованные в проекте", "14")
    add_card(slide, Inches(0.72), Inches(2.1), Inches(5.5), Inches(1.45), "Аутентификация и доступ", "JWT-токен используется для входа в админ-панель, а CRUD-операции над контентом и заказами защищены зависимостью get_current_admin_user.", body_size=14)
    add_card(slide, Inches(6.4), Inches(2.1), Inches(5.9), Inches(1.45), "Валидация и защита ввода", "На клиенте и сервере выполняется проверка данных форм, ограничение типа загружаемых файлов и обработка ошибок через HTTPException.", body_size=14, fill=PANEL_ALT)
    add_card(slide, Inches(0.72), Inches(3.9), Inches(5.5), Inches(1.55), "Хранение и учетные данные", "Пароли администратора хешируются через bcrypt/passlib, чувствительные параметры вынесены в .env, а БД хранится отдельно от frontend-контейнера.", body_size=14)
    add_card(slide, Inches(6.4), Inches(3.9), Inches(5.9), Inches(1.55), "Инфраструктурная устойчивость", "Docker-изоляция разделяет frontend и backend, каталоги images и restaurant.db вынесены в постоянное хранилище, что упрощает резервное копирование.", body_size=14, fill=PANEL_ALT)
    add_footer(slide, 14)

    # 15
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "Экономика проекта", "Оценка трудозатрат и ожидаемого эффекта от внедрения", "15")
    add_card(slide, Inches(0.72), Inches(2.0), Inches(6.55), Inches(4.6), "Оценка трудозатрат", "1. Анализ требований и структура — 24 ч × 500 ₽ = 12 000 ₽\n2. Frontend (Vue SPA) — 72 ч × 550 ₽ = 39 600 ₽\n3. Backend и БД — 64 ч × 550 ₽ = 35 200 ₽\n4. Тестирование и исправления — 28 ч × 500 ₽ = 14 000 ₽\n5. Docker, документация и презентация — 24 ч × 500 ₽ = 12 000 ₽\n\nИтого ориентировочная стоимость разработки: 112 800 ₽", body_size=15, fit_body=True, body_top=0.6)
    add_metric(slide, Inches(8.18), Inches(2.15), Inches(3.65), Inches(1.18), "≈ 113 тыс. ₽", "оценка стоимости разработки", number_size=20, label_size=11)
    add_metric(slide, Inches(8.18), Inches(3.68), Inches(3.65), Inches(1.18), "≈ 25–35 тыс. ₽/мес", "потенциальный ежемесячный эффект", number_size=18, label_size=11)
    add_metric(slide, Inches(8.18), Inches(5.21), Inches(3.65), Inches(1.18), "4–5 месяцев", "ориентировочный срок окупаемости", number_size=19, label_size=11)
    add_footer(slide, 15)

    # 16
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_title(slide, "QR-код на программный продукт", "Демо-ссылка для быстрого перехода к развернутому приложению", "16")
    if OUTPUT_QR.exists():
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(2.12), Inches(3.9), Inches(3.9))
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(250, 250, 250)
        frame.line.fill.background()
        slide.shapes.add_picture(str(OUTPUT_QR), Inches(1.33), Inches(2.45), height=Inches(3.2))
    add_card(
        slide,
        Inches(5.85),
        Inches(2.12),
        Inches(6.0),
        Inches(4.18),
        "Демо-стенд",
        f"URL для открытия проекта:\n{qr_url}\n\nЧто доступно по ссылке:\n• публичный frontend ресторана\n• проксирование API через nginx\n• backend FastAPI в Docker\n\nДанные для входа в админку:\nadmin / admin123",
        fill=PANEL_ALT,
        body_size=16,
        fit_body=True,
        body_top=0.64,
    )
    add_footer(slide, 16)

    # 17
    slide = prs.slides.add_slide(blank)
    add_bg(slide, bg_variant)
    add_textbox(slide, Inches(0.75), Inches(1.35), Inches(8.8), Inches(0.4), "СПАСИБО ЗА ВНИМАНИЕ", font_size=14, bold=True, color=ACCENT)
    add_textbox(slide, Inches(0.72), Inches(1.85), Inches(8.8), Inches(1.1), "Разработка веб-приложения\nдля ресторана «Гастрономия»", font_size=30, bold=True)
    add_textbox(slide, Inches(0.75), Inches(3.2), Inches(8.5), Inches(0.45), "Проект объединяет клиентский сайт, формы заказа/бронирования, админ-панель и Docker-развертывание.", font_size=16, color=TEXT_MUTED)
    add_card(slide, Inches(0.75), Inches(5.45), Inches(3.55), Inches(1.08), "Студент", "Давранбеков С.Б.", fit_body=True)
    add_card(slide, Inches(4.55), Inches(5.45), Inches(3.55), Inches(1.08), "Проект", "Ресторан «Гастрономия»", fit_body=True)
    add_card(slide, Inches(8.35), Inches(5.45), Inches(3.75), Inches(1.08), "Демо", qr_url, fill=PANEL_ALT, fit_body=True)

    prs.save(str(output_path))


if __name__ == "__main__":
    build_presentation(OUTPUT_PPTX, bg_variant="clean")
    build_presentation(OUTPUT_PPTX_CLEAN, bg_variant="clean")
    build_presentation(OUTPUT_PPTX_SOFT, bg_variant="soft")
